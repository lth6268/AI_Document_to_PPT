
"""
app.py — 文档转 PPT 生成器（Flask + 本地 Ollama）
支持 .docx / .pdf 上传，按页生成 PPTX，图片保留在底部。
"""

import os
import io
import json
import uuid
import re
from pathlib import Path

from flask import Flask, request, jsonify, send_file, render_template
from openai import OpenAI
from docx import Document
from docx.oxml.ns import qn
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ============================================================
# 配置
# ============================================================
app = Flask(__name__)
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), "uploads")
OUTPUT_FOLDER = os.path.join(os.path.dirname(__file__), "outputs")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Ollama 客户端（兼容 OpenAI 接口）
OLLAMA_BASE = os.environ.get("OLLAMA_BASE", "http://localhost:11434/v1")
OLLAMA_API_KEY = os.environ.get("OLLAMA_API_KEY", "ollama")
# 后端默认模型设置处 —— 修改这里即可更改默认模型
DEFAULT_MODEL = os.environ.get("DEFAULT_MODEL", "deepseek-r1:8b")

client = OpenAI(base_url=OLLAMA_BASE, api_key=OLLAMA_API_KEY)

# PPT 尺寸 (16:9 横向)
PPT_WIDTH = Inches(13.33)
PPT_HEIGHT = Inches(7.5)
TEXT_AREA_TOP = Inches(0.5)
TEXT_AREA_HEIGHT = Inches(4.5)
IMG_AREA_TOP = Inches(5.3)
IMG_AREA_HEIGHT = Inches(1.8)

# 颜色
COLOR_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x7A, 0xCC)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_TABLE_HEADER = RGBColor(0x2D, 0x5F, 0x8A)
COLOR_TABLE_CELL = RGBColor(0xE8, 0xEE, 0xF4)


# ============================================================
# 文档解析
# ============================================================

def parse_docx(file_path):
    """解析 Word 文档，按段落逻辑分页，提取图片"""
    doc = Document(file_path)
    pages_data = []
    current_page_text = []
    current_page_images = []
    char_count = 0
    PAGE_LIMIT = 800  # 约 800 字一页

    def save_page():
        nonlocal current_page_text, current_page_images, char_count
        if current_page_text or current_page_images:
            pages_data.append({
                "text": "\n".join(current_page_text),
                "images": current_page_images[:]
            })
            current_page_text = []
            current_page_images = []
            char_count = 0

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        # 检测分页符
        if tag == "p":
            for run in element.iter(qn("w:r")):
                br = run.find(qn("w:br"))
                if br is not None and br.get(qn("w:type")) == "page":
                    save_page()

            # 段落文本
            full_text = ""
            for t in element.iter(qn("w:t")):
                if t.text:
                    full_text += t.text
            full_text = full_text.strip()
            if full_text:
                current_page_text.append(full_text)
                char_count += len(full_text)
                if char_count > PAGE_LIMIT:
                    save_page()

            # 段落中的图片
            for drawing in element.iter(qn("w:drawing")):
                for blip in drawing.iter(qn("a:blip")):
                    rId = blip.get(qn("r:embed"))
                    if rId:
                        try:
                            image_part = doc.part.related_parts[rId]
                            img_bytes = image_part.blob
                            ct = image_part.content_type
                            ext = ct.split("/")[-1]
                            if ext == "jpeg":
                                ext = "jpg"
                            img_name = f"{uuid.uuid4()}.{ext}"
                            img_path = os.path.join(UPLOAD_FOLDER, img_name)
                            with open(img_path, "wb") as f:
                                f.write(img_bytes)
                            current_page_images.append(img_path)
                        except Exception as e:
                            print(f"  [docx图片提取失败] {e}")

        # 表格中的图片
        if tag == "tbl":
            for drawing in element.iter(qn("w:drawing")):
                for blip in drawing.iter(qn("a:blip")):
                    rId = blip.get(qn("r:embed"))
                    if rId:
                        try:
                            image_part = doc.part.related_parts[rId]
                            img_bytes = image_part.blob
                            ct = image_part.content_type
                            ext = ct.split("/")[-1]
                            if ext == "jpeg":
                                ext = "jpg"
                            img_name = f"{uuid.uuid4()}.{ext}"
                            img_path = os.path.join(UPLOAD_FOLDER, img_name)
                            with open(img_path, "wb") as f:
                                f.write(img_bytes)
                            current_page_images.append(img_path)
                        except Exception as e:
                            print(f"  [docx表格图片提取失败] {e}")

    save_page()
    return pages_data


def parse_pdf(file_path):
    """解析 PDF，严格 1:1 分页，提取图片和文本"""
    pages_data = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text() or ""
            page_images = []

            # 提取页面中的图片
            images = page.images
            for i, img_obj in enumerate(images):
                try:
                    crop_box = (img_obj["x0"], img_obj["top"], img_obj["x1"], img_obj["bottom"])
                    cropped = page.crop(crop_box)
                    img_path = os.path.join(
                        UPLOAD_FOLDER, f"p{len(pages_data)}_img{i}.png"
                    )
                    cropped.to_image(resolution=200).save(img_path)
                    page_images.append(img_path)
                except Exception as e:
                    print(f"  [pdf图片提取失败] {e}")

            pages_data.append({
                "text": page_text,
                "images": page_images
            })
    return pages_data


def parse_document(file_path):
    ext = Path(file_path).suffix.lower()
    if ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    else:
        raise ValueError(f"不支持的格式: {ext}")


# ============================================================
# 调用 Ollama 提炼内容
# ============================================================

def summarize_with_ollama(text, model):
    """调用本地大模型提炼标题和要点"""
    if not text.strip():
        return {"title": "无内容", "points": ["页面内容为空或仅为图片"]}

    truncated = text[:4000]
    prompt = f"""你是一位PPT内容提炼助手。请将以下文本精简为适合PPT展示的内容。
            要求：
            1. 优先使用原文语言，如原文语言为俄语则使用俄语，原文语言为中文则使用中文。
            2. 标题不超过15字，直接点明核心主题，不要修饰词。
            3. 提炼3-6个要点，每条不超过100字，保留最核心的信息，去掉所有修饰和解释。
            4. 不要使用'首先''其次''此外'等过渡词，直接陈述要点。
            5. 必须输出严格的JSON格式，不要输出其他任何内容，格式如下：
            {{
                "title": "这里填标题",
                "content": ["要点1", "要点2", "要点3"]
            }}
            不要包含markdown标记（如 ```json），不要包含任何解释性文字。
            
            原文内容：
            {truncated}"""
    
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3,
        )
        content = response.choices[0].message.content.strip()
        # 兼容 markdown 代码块
        json_match = re.search(r"```(?:json)?\s*(\{.*?\})\s*```", content, re.DOTALL)
        if json_match:
            content = json_match.group(1)
        return json.loads(content)
    except Exception as e:
        print(f"  [Ollama错误] {e}")
        return {"title": "处理出错", "points": [str(e)[:80]]}


# ============================================================
# 生成 PPTX
# ============================================================

def create_ppt(pages_data, outline_title, model, output_path):
    """生成 PPTX，使用传入的 model 调用 Ollama"""
    prs = Presentation()
    prs.slide_width = PPT_WIDTH
    prs.slide_height = PPT_HEIGHT

    blank_layout = prs.slide_layouts[6]  # 空白布局

    for idx, page_info in enumerate(pages_data):
        slide = prs.slides.add_slide(blank_layout)

        # 设置背景色
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # 1. 调用 AI 总结（使用前端传来的 model）
        ai_result = summarize_with_ollama(page_info["text"], model)
        title = ai_result.get("title", "无标题")
        points = ai_result.get("content", [])  # 注意：这里是 content 不是 points

        # 2. 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT
        p.alignment = PP_ALIGN.LEFT

        # 标题下装饰线
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.2), Inches(12), Pt(3)  # 从1.15调整为1.2
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_ACCENT
        line.line.fill.background()

        # 3. 添加正文要点（调整起始位置，避免与标题重叠）
        BODY_TOP = Inches(1.4)  # 原来0.5，现在改为1.4，在标题下方
        BODY_HEIGHT = Inches(3.8)  # 原来4.5，调整为3.8
        
        body_box = slide.shapes.add_textbox(
            Inches(0.6), 
            BODY_TOP, 
            Inches(12), 
            BODY_HEIGHT
        )
        tf_body = body_box.text_frame
        tf_body.word_wrap = True

        # 如果 points 为空，显示提示信息
        if not points:
            p = tf_body.paragraphs[0]
            p.text = "  （无文本内容）"
            p.font.size = Pt(18)
            p.font.color.rgb = COLOR_TEXT
        else:
            for i, point in enumerate(points):
                p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
                # 确保 point 是字符串
                point_text = str(point) if point else "（空内容）"
                p.text = f"  {point_text}"
                p.font.size = Pt(18)
                p.font.color.rgb = COLOR_TEXT
                p.space_after = Pt(12)

        # 4. 添加图片（底部横向排列）
        images = page_info.get("images", [])
        if images:
            # 图片区域背景（位置不变）
            img_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), IMG_AREA_TOP, PPT_WIDTH, IMG_AREA_HEIGHT
            )
            img_bg.fill.solid()
            img_bg.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF4)
            img_bg.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            img_bg.line.width = Pt(1)

            # 图片标签
            label_box = slide.shapes.add_textbox(Inches(0.3), IMG_AREA_TOP + Inches(0.05), Inches(3), Inches(0.35))
            tf_label = label_box.text_frame
            p = tf_label.paragraphs[0]
            p.text = "📷 原文图片"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
            p.font.bold = True

            # 排列图片
            max_imgs = min(len(images), 3)
            display_imgs = images[:max_imgs]
            total_gap = Inches(0.3)
            single_w = (PPT_WIDTH - Inches(0.6) - total_gap * (max_imgs - 1)) / max_imgs
            single_h = IMG_AREA_HEIGHT - Inches(0.3)

            for j, img_path in enumerate(display_imgs):
                try:
                    left_pos = Inches(0.3) + j * (single_w + total_gap)
                    top_pos = IMG_AREA_TOP + Inches(0.25)
                    slide.shapes.add_picture(img_path, left_pos, top_pos, width=single_w, height=single_h)
                except Exception as e:
                    print(f"  [插入图片失败] {img_path}: {e}")

    prs.save(output_path)
    return output_path

# ============================================================
# Web 路由
# ============================================================

@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/models")
def get_models():
    """获取 Ollama 已安装的模型列表，并返回后端默认模型"""
    try:
        import urllib.request
        resp = urllib.request.urlopen("http://localhost:11434/api/tags", timeout=5)
        data = json.loads(resp.read())
        models = [m["name"] for m in data.get("models", [])]
        return jsonify({
            "models": models,
            "default_model": DEFAULT_MODEL
        })
    except Exception as e:
        return jsonify({"models": [], "default_model": DEFAULT_MODEL, "error": str(e)})


@app.route("/api/generate", methods=["POST"])
def generate():
    """上传文档 → 解析 → Ollama提炼 → 生成PPTX → 返回下载链接"""
    file = request.files.get("file")
    # 从前端表单读取用户选择的模型
    model = request.form.get("model", DEFAULT_MODEL)

    if not file:
        return jsonify({"error": "请上传文件"}), 400

    filename = file.filename
    ext = Path(filename).suffix.lower()
    if ext not in (".docx", ".pdf"):
        return jsonify({"error": "不支持的格式，请上传 .docx 或 .pdf"}), 400

    # 保存上传文件
    upload_path = os.path.join(UPLOAD_FOLDER, f"{uuid.uuid4()}{ext}")
    file.save(upload_path)

    try:
        # 1. 解析文档
        print(f"[解析] {filename} (格式: {ext})")
        pages_data = parse_document(upload_path)
        print(f"  解析得到 {len(pages_data)} 页")

        if not pages_data:
            return jsonify({"error": "文档内容为空"}), 400

        # 2. 生成 PPT（使用前端传来的 model）
        first_title = "演示文稿"
        try:
            ai_first = summarize_with_ollama(pages_data[0]["text"], model)
            first_title = ai_first.get("title", "演示文稿")
        except:
            pass

        output_filename = f"{uuid.uuid4().hex[:8]}.pptx"
        output_path = os.path.join(OUTPUT_FOLDER, output_filename)
        # 将 model 传递给 create_ppt
        create_ppt(pages_data, first_title, model, output_path)

        # 清理上传文件
        os.remove(upload_path)

        return jsonify({
            "success": True,
            "download": f"/download/{output_filename}",
            "slides_count": len(pages_data),
            "title": first_title,
        })

    except Exception as e:
        import traceback
        traceback.print_exc()
        # 清理残留文件
        try:
            os.remove(upload_path)
        except:
            pass
        return jsonify({"error": str(e)}), 500


@app.route("/download/<filename>")
def download(filename):
    filepath = os.path.join(OUTPUT_FOLDER, filename)
    if not os.path.exists(filepath):
        return jsonify({"error": "文件不存在"}), 404
    return send_file(filepath, as_attachment=True, download_name=filename)


# 全局错误处理
@app.errorhandler(500)
def internal_error(error):
    return jsonify({"error": f"服务器内部错误: {str(error)}"}), 500


@app.errorhandler(404)
def not_found(error):
    return jsonify({"error": "接口不存在"}), 404


if __name__ == "__main__":
    print(f"🚀 启动服务: http://localhost:6001")
    print(f"   Ollama API: {OLLAMA_BASE}")
    print(f"   默认模型: {DEFAULT_MODEL}")
    app.run(debug=False, host="0.0.0.0", port=6001)
